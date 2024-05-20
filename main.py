from pptx import Presentation
from PIL import Image

def convert_ppt_to_images(ppt_file_path, output_folder):
    # Load the PowerPoint presentation
    presentation = Presentation(ppt_file_path)

    # Iterate through each slide in the presentation
    for i, slide in enumerate(presentation.slides):
        # Create a new image with white background
        img = Image.new('RGB', (int(presentation.slide_width), int(presentation.slide_height)), 'white')

        # Create a drawing object to render the slide onto the image
        draw = ImageDraw.Draw(img)

        # Render each shape on the slide onto the image
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                # If the shape is an image, paste it onto the image
                image = Image.open(io.BytesIO(shape.image.blob))
                img.paste(image, (int(shape.left), int(shape.top)))

            if shape.has_text_frame:
                # If the shape has text, draw it onto the image
                text = shape.text
                font = ImageFont.load_default()
                draw.text((int(shape.left), int(shape.top)), text, font=font, fill='black')

        # Save the image to the output folder
        img.save(f'{output_folder}/slide_{i + 1}.png')

if __name__ == "__main__":
    ppt_file_path = "C:\\Users\\hp\\OneDrive\\Documents\\Presentation_on_Jini.pptx"
    output_folder = "C:\\Users\\hp\\OneDrive\\Desktop\\New folder"

    convert_ppt_to_images(ppt_file_path, output_folder)
